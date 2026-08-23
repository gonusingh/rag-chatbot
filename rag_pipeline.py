# ============================================================
#  RAG PIPELINE v5 - rag_pipeline.py
#  SUPPORTS: PDF, TXT, DOCX, XLSX, PPTX, Web URLs
#  NEW IN v5: Proper logging throughout
#
#  LOGGING LEVELS (in order of severity):
#  DEBUG    → very detailed, for development only
#  INFO     → normal operations ("file loaded", "chunks created")
#  WARNING  → something unexpected but not breaking
#  ERROR    → something failed
#  CRITICAL → entire app is broken
# ============================================================


# ── IMPORTS ─────────────────────────────────────────────────
import logging
import os
import numpy as np
import faiss
import fitz
from sentence_transformers import SentenceTransformer
from dotenv import load_dotenv
import re

load_dotenv()


# ── LOGGING SETUP ────────────────────────────────────────────
# Must be configured BEFORE any logger.info() calls
# Otherwise messages get dropped silently

logging.basicConfig(
    level=logging.INFO,
    # Show INFO and above (INFO, WARNING, ERROR, CRITICAL)
    # Change to logging.DEBUG to see everything

    format="%(asctime)s %(levelname)-8s %(message)s",
    # Example output:
    # 2024-01-15 14:23:01 INFO     Reading PDF: resume.pdf
    # 2024-01-15 14:23:05 WARNING  Low relevance: 0.45
    # 2024-01-15 14:23:06 ERROR    URL timeout: example.com

    datefmt="%Y-%m-%d %H:%M:%S",
    # Date format: Year-Month-Day Hour:Minute:Second

    handlers=[
        logging.StreamHandler(),
        # StreamHandler → prints to terminal

        logging.FileHandler("rag_app.log", encoding='utf-8')
        # FileHandler → saves to rag_app.log file
        # 'a' mode (default) = append, never overwrites
        # Every run adds to the log file
    ]
) 

# Create module-level logger
# Best practice: one logger per file, named after the module
logger = logging.getLogger("RAG_Pipeline")


# ============================================================
#  STEP 1: DOCUMENT LOADERS
#
#  SUPPORTED FORMATS:
#  .pdf  → PyMuPDF (fitz)
#  .txt  → built-in open()
#  .docx → python-docx
#  .xlsx → openpyxl
#  .pptx → python-pptx
#  URLs  → requests + BeautifulSoup
# ============================================================

def load_pdf(file_path):
    """
    Reads every page of a PDF and returns all text as one string.

    PDF stores text as graphical elements with positions/fonts.
    PyMuPDF (fitz) extracts just the readable text content.
    """
    logger.info(f"Reading PDF: {file_path}")

    if not os.path.exists(file_path):
        logger.error(f"PDF not found: {file_path}")
        raise FileNotFoundError(f"PDF not found: {file_path}")

    pdf_document = fitz.open(file_path)
    all_text = []

    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        page_text = page.get_text()
        all_text.append(page_text)
        logger.debug(f"Page {page_num + 1}: {len(page_text)} chars")

    page_count = len(pdf_document)
    pdf_document.close()
    full_text = "\n".join(all_text)

    logger.info(f"PDF loaded: {len(full_text)} characters, "
                f"{page_count} pages")
    return full_text


def load_txt(file_path):
    """
    Reads a plain .txt file and returns its content.
    Simplest loader — text files are just raw text.
    """
    logger.info(f"Reading TXT: {file_path}")

    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()

    logger.info(f"TXT loaded: {len(text)} characters")
    return text


def load_docx(file_path):
    """
    Reads a Word .docx file and returns all paragraph text.

    Word document structure:
    Document → Paragraphs → Runs → Text

    We grab text from every paragraph and join them.
    Empty paragraphs (blank lines) are skipped.
    """
    logger.info(f"Reading DOCX: {file_path}")

    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    from docx import Document

    doc = Document(file_path)
    all_text = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            all_text.append(paragraph.text)

    full_text = "\n".join(all_text)

    logger.info(f"DOCX loaded: {len(doc.paragraphs)} paragraphs, "
                f"{len(full_text)} characters")
    return full_text


def load_xlsx(file_path):
    """
    Reads an Excel .xlsx file and converts rows to readable text.

    Excel structure:
    Workbook → Sheets → Rows → Cells

    Converts each row to:
    "Header1: Value1 | Header2: Value2 | Header3: Value3"

    Example:
    Name: Vinit | Skills: Python | Experience: 2 years
    """
    logger.info(f"Reading Excel: {file_path}")

    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    import openpyxl

    # data_only=True reads VALUES not formulas
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    all_text = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        logger.debug(f"Reading sheet: '{sheet_name}'")
        all_text.append(f"\n=== Sheet: {sheet_name} ===")

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            logger.warning(f"Sheet '{sheet_name}' is empty")
            continue

        # First row = headers
        headers = [
            str(h) if h is not None else f"Column_{i}"
            for i, h in enumerate(rows[0])
        ]

        # Remaining rows = data
        row_count = 0
        for row in rows[1:]:
            if all(cell is None for cell in row):
                continue
            row_parts = []
            for header, value in zip(headers, row):
                if value is not None:
                    row_parts.append(f"{header}: {value}")
            if row_parts:
                all_text.append(" | ".join(row_parts))
                row_count += 1

        logger.debug(f"Sheet '{sheet_name}': {row_count} rows extracted")

    full_text = "\n".join(all_text)
    logger.info(f"Excel loaded: {len(workbook.sheetnames)} sheets, "
                f"{len(full_text)} characters")
    return full_text


def load_pptx(file_path):
    """
    Reads a PowerPoint .pptx file and extracts text from all slides.

    PPTX structure:
    Presentation → Slides → Shapes → TextFrame → Paragraphs → Text

    hasattr(shape, 'text_frame') check is important:
    Images, charts, icons don't have text frames.
    Only textboxes and titles do.
    """
    logger.info(f"Reading PowerPoint: {file_path}")

    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    from pptx import Presentation

    prs = Presentation(file_path)
    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        all_text.append(f"\n=== Slide {slide_num} ===")
        slide_text_count = 0

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        all_text.append(text)
                        slide_text_count += 1

        logger.debug(f"Slide {slide_num}: {slide_text_count} text blocks")

    full_text = "\n".join(all_text)
    logger.info(f"PPTX loaded: {len(prs.slides)} slides, "
                f"{len(full_text)} characters")
    return full_text


def load_url(url):
    """
    Fetches a webpage and extracts clean readable text.

    HOW WEB SCRAPING WORKS:
    1. requests.get(url) downloads raw HTML
    2. BeautifulSoup parses HTML into a tree
    3. We remove noise: scripts, styles, nav, footer
    4. Extract only meaningful text content

    Returns tuple: (text, url) so source can be shown in UI.
    """
    import requests
    from bs4 import BeautifulSoup

    logger.info(f"Fetching URL: {url}")

    # Validate URL format
    if not url.startswith(('http://', 'https://')):
        logger.error(f"Invalid URL format: {url}")
        raise ValueError(
            f"Invalid URL: '{url}'\n"
            f"Must start with http:// or https://"
        )

    # Headers to pretend to be a real browser
    # Some websites block requests without User-Agent
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/91.0.4472.124 Safari/537.36"
        )
    }

    try:
        response = requests.get(
            url, headers=headers, timeout=10
        )
        # Raises exception for 4xx/5xx status codes
        response.raise_for_status()
        logger.info(f"URL fetched: status {response.status_code}")

    except requests.exceptions.ConnectionError:
        logger.error(f"Cannot connect to URL: {url}")
        raise ConnectionError(
            f"Cannot reach: {url}\n"
            f"Check internet connection or URL."
        )
    except requests.exceptions.Timeout:
        logger.error(f"URL timed out: {url}")
        raise TimeoutError(
            f"Timed out after 10 seconds: {url}"
        )
    except requests.exceptions.HTTPError as e:
        logger.error(f"HTTP error for {url}: {e}")
        raise Exception(
            f"HTTP Error for {url}: {str(e)}"
        )

    # Parse HTML with BeautifulSoup
    soup = BeautifulSoup(response.text, 'html.parser')

    # Remove noise tags
    noise_tags = [
        'script', 'style', 'nav', 'footer',
        'header', 'aside', 'advertisement'
    ]
    removed = 0
    for tag in soup.find_all(noise_tags):
        tag.decompose()
        removed += 1
    logger.debug(f"Removed {removed} noise tags")

    # Try to find main content area
    main_content = (
        soup.find('main') or
        soup.find('article') or
        soup.find('div', {'id': 'content'}) or
        soup.find('div', {'class': 'content'}) or
        soup.body
    )

    if main_content:
        text = main_content.get_text(separator='\n', strip=True)
    else:
        logger.warning(f"No main content found for: {url}")
        text = soup.get_text(separator='\n', strip=True)

    # Clean up text
    lines = [line.strip() for line in text.splitlines()]
    lines = [line for line in lines if line]
    lines = [line for line in lines if len(line) > 20]
    clean_text = '\n'.join(lines)

    logger.info(f"URL loaded: {len(clean_text)} characters extracted")
    return clean_text, url


def load_document(file_path):
    """
    MASTER LOADER — routes to correct loader by file extension.

    .pdf  → load_pdf()     PyMuPDF
    .txt  → load_txt()     built-in
    .docx → load_docx()    python-docx
    .xlsx → load_xlsx()    openpyxl
    .pptx → load_pptx()    python-pptx
    URLs  → load_url()     requests + BeautifulSoup
    """
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()
    logger.info(f"Loading document: type='{extension}'")

    if extension == '.pdf':
        return load_pdf(file_path)
    elif extension == '.txt':
        return load_txt(file_path)
    elif extension == '.docx':
        return load_docx(file_path)
    elif extension == '.xlsx':
        return load_xlsx(file_path)
    elif extension == '.pptx':
        return load_pptx(file_path)
    else:
        logger.error(f"Unsupported file type: '{extension}'")
        raise ValueError(
            f"Unsupported file: '{extension}'\n"
            f"Supported: .pdf, .txt, .docx, .xlsx, .pptx"
        )


# ============================================================
#  STEP 2: CHUNKING
#
#  Splits text into overlapping pieces for precise retrieval.
#
#  KEY FORMULA: start += chunk_size - overlap
#
#  WHY OVERLAP?
#  Sentences near boundaries get split without overlap.
#  Overlap ensures boundary sentences appear complete
#  in at least one chunk → better retrieval quality.
#
#  SIZE GUIDE:
#  1-2 pages    → 100-150 words, overlap 20-30
#  5-20 pages   → 200-300 words, overlap 50
#  50-100 pages → 300-400 words, overlap 75-100
# ============================================================

def chunk_text(text, chunk_size=200, overlap=50):
    """
    Splits text into overlapping chunks.

    Args:
        text: full document text
        chunk_size: words per chunk
        overlap: words repeated between chunks

    Returns:
        list of chunk strings
    """
    logger.info(f"Chunking text: size={chunk_size}, overlap={overlap}")

    words = text.split()
    logger.info(f"Total words: {len(words)}")

    if len(words) == 0:
        logger.warning("Empty text provided to chunk_text!")
        return []

    chunks = []
    start = 0

    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap

    logger.info(f"Created {len(chunks)} chunks")

    if chunks:
        logger.debug(f"Chunk preview: '{chunks[0][:80]}...'")

    return chunks


# ============================================================
#  STEP 3: EMBEDDINGS
#
#  Converts text → 384 numbers (vector).
#  Similar meaning → similar numbers → small L2 distance.
#
#  WHY all-MiniLM-L6-v2?
#  - Small (~90MB), fast, free, runs locally
#  - Good quality for English text
#  - 384 dimensions — good balance of speed vs quality
# ============================================================

def load_embedding_model():
    """
    Loads the sentence transformer model.
    First run downloads ~90MB (cached forever after).
    """
    logger.info("Loading embedding model: all-MiniLM-L6-v2")
    model = SentenceTransformer('all-MiniLM-L6-v2')
    logger.info("Embedding model loaded successfully")
    return model


def create_embeddings(chunks, model):
    """
    Converts text chunks → numpy array of shape (n_chunks, 384).

    INPUT:  ["Photosynthesis is...", "Water cycle..."]
    OUTPUT: numpy array shape (2, 384)

    float32 required by FAISS.
    """
    logger.info(f"Creating embeddings for {len(chunks)} chunks")
    embeddings = model.encode(chunks, show_progress_bar=True)
    embeddings = np.array(embeddings, dtype=np.float32)
    logger.info(f"Embeddings created: shape={embeddings.shape}")
    return embeddings


# ============================================================
#  STEP 4: FAISS INDEX
#
#  Stores all vectors for fast similarity search.
#  IndexFlatL2 = exact search using Euclidean distance.
#  Lower distance = more similar = better match.
# ============================================================

def build_faiss_index(embeddings):
    """
    Stores embeddings in FAISS for fast similarity search.

    Args:
        embeddings: numpy array shape (n_chunks, 384)

    Returns:
        faiss.IndexFlatL2: searchable vector index
    """
    logger.info("Building FAISS index...")
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings)
    logger.info(f"FAISS index built: {index.ntotal} vectors stored")
    return index


# ============================================================
#  STEP 5: RETRIEVAL
#
#  Finds most relevant chunks for user's question.
#  MUST use same embedding model as chunks!
#  Lower distance = higher relevance.
# ============================================================

def retrieve_chunks(query, model, index, chunks, top_k=3):
    """
    Finds top-k most relevant chunks for the query.

    Args:
        query: user's question string
        model: same SentenceTransformer used for chunks
        index: FAISS index containing chunk vectors
        chunks: original text chunk strings
        top_k: number of chunks to retrieve

    Returns:
        list of top-k most relevant chunk strings
    """
    logger.info(f"Retrieving chunks for query: '{query[:50]}...'")

    query_vector = model.encode([query])
    query_vector = np.array(query_vector, dtype=np.float32)

     # Never request more chunks than we actually have
    safe_top_k = min(top_k, len(chunks))
    if safe_top_k < top_k:
     logger.warning(
        f"Requested top_k={top_k} but only {len(chunks)} "
        f"chunks available. Using top_k={safe_top_k}"
    )
    distances, indices = index.search(query_vector, safe_top_k)

    retrieved = []
    for rank, (idx, dist) in enumerate(
        zip(indices[0], distances[0])
    ):
         # Guard against overflow with huge distances
        try:
           relevance = max(0, 100 - (dist * 20))
        except (OverflowError, ValueError):
           relevance = 0
        logger.info(
            f"Match #{rank+1}: chunk={idx}, "
            f"distance={dist:.3f}, relevance={relevance:.0f}%"
        )

        if relevance < 30:
            logger.warning(
                f"Low relevance score ({relevance:.0f}%) "
                f"for chunk #{idx} — answer may be inaccurate"
            )

        retrieved.append(chunks[idx])

    return retrieved


# ============================================================
#  STEP 6: GENERATION
#
#  Sends retrieved chunks + question to Groq (Llama 3).
#
#  PROMPT STRUCTURE:
#  Role + Context (chunks) + Question + Rules → Answer
#
#  "ONLY use context" prevents hallucination.
#  Temperature=0.1 = focused/factual (not creative).
# ============================================================

def ask_llm(query, retrieved_chunks, api_key):
    """
    Sends context + question to Groq and returns answer.

    Args:
        query: user's question
        retrieved_chunks: top-k relevant text chunks
        api_key: Groq API key from .env

    Returns:
        answer string from Llama 3
    """
    try:
        from groq import Groq
    except ImportError:
        logger.error("Groq library not installed!")
        return "❌ Run: pip install groq"

    logger.info("Sending request to Groq (Llama 3)...")

    client = Groq(api_key=api_key)
    context = "\n\n--- (next chunk) ---\n\n".join(retrieved_chunks)

    prompt = f"""You are a helpful assistant. Answer the question using ONLY the context below.

CONTEXT FROM DOCUMENT:
{context}

QUESTION: {query}

RULES:
- Use only information from the context above
- If the answer is not in context, say "This information is not in the provided document"
- Be clear and concise

ANSWER:"""

    try:
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1
        )
        answer = response.choices[0].message.content
        answer = re.sub(
            r'<think>.*?</think>', '',
            answer,
            flags=re.DOTALL
        ).strip()
        logger.info(
            f"Groq response received: {len(answer)} characters"
        )
        return answer

    except Exception as e:
        logger.error(f"Groq API error: {str(e)}")
        raise


# ============================================================
#  MAIN — Terminal version for quick testing
# ============================================================

def main():
    logger.info("=" * 50)
    logger.info("RAG PIPELINE v5 STARTING")
    logger.info("=" * 50)

    USE_PDF    = False
    PDF_PATH   = "sample.pdf"
    CHUNK_SIZE = 80
    OVERLAP    = 20
    TOP_K      = 3

    if USE_PDF:
        text = load_document(PDF_PATH)
    else:
        logger.info("Using built-in sample text")
        text = """
        Photosynthesis is the process by which plants use sunlight,
        water, and carbon dioxide to produce oxygen and energy in the
        form of sugar. This process takes place in the chloroplasts,
        using the green pigment called chlorophyll.

        The water cycle describes how water evaporates from the surface
        of the earth, rises into the atmosphere, cools and condenses
        into rain or snow in clouds, and falls again to the surface.

        Gravity is a fundamental force that attracts objects with mass
        toward each other. Isaac Newton described it mathematically,
        while Einstein explained it as curvature of spacetime.

        The digestive system breaks down food into nutrients. It starts
        in the mouth, moves to the stomach with gastric acids, then
        the small intestine absorbs nutrients and large intestine
        absorbs water.
        """

    model      = load_embedding_model()
    chunks     = chunk_text(text, chunk_size=CHUNK_SIZE, overlap=OVERLAP)
    embeddings = create_embeddings(chunks, model)
    index      = build_faiss_index(embeddings)

    api_key = os.getenv("GROQ_API_KEY", "")
    if not api_key:
        logger.warning("No GROQ_API_KEY found in .env file")
    else:
        logger.info("Groq API key loaded")

    logger.info("RAG Pipeline ready! Type 'quit' to exit.")

    while True:
        query = input("\n❓ Your question: ").strip()

        if query.lower() in ['quit', 'exit', 'q']:
            logger.info("User exited the application")
            print("👋 Goodbye!")
            break

        if not query:
            continue

        relevant_chunks = retrieve_chunks(
            query, model, index, chunks, top_k=TOP_K
        )

        print("\n📄 Retrieved Context:")
        for i, chunk in enumerate(relevant_chunks, 1):
            print(f"\n  [{i}] {chunk[:120]}...")

        if api_key:
            answer = ask_llm(query, relevant_chunks, api_key)
            print(f"\n💬 Answer:\n{answer}")
        else:
            print("\n💡 Add GROQ_API_KEY to .env for AI answers")


if __name__ == "__main__":
    main()